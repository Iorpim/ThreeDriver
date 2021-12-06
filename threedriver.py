#!/usr/env python3
from selenium import webdriver
from selenium.webdriver.common import action_chains
from selenium.webdriver.common.proxy import Proxy, ProxyType
from selenium.webdriver.common.keys import Keys
from selenium.webdriver.common.action_chains import ActionChains
import time
import argparse
import json
import traceback
import os
import re
import humanfriendly
import docx
import pptx
import openpyxl
import magic
import PyPDF2
import extract_msg

DOWNLOAD_FOLDER = ""
URL = "https://onedrive.live.com/about/en-us/signin/"
PROXY = False
FIREFOX = False
def __init__():
    global DOWNLOAD_FOLDER
    global URL
    global driver
    DOWNLOAD_FOLDER = os.path.join(os.getcwd(), "Downloads_temp")
    options = webdriver.ChromeOptions()
    if(PROXY):
        options.add_argument(f"--proxy-server={PROXY}")
    #options.add_argument("--proxy-server=http://127.0.0.1:8080")
    try:
        os.mkdir("Downloads_temp")
    except FileExistsError:
        try:
            os.rename("Downloads_temp", "Downloads_temp_backup")
        except FileExistsError: # TODO Deal better with this
            os.removedirs("Downloads_temp")
        os.mkdir("Downloads_temp")
    options.add_experimental_option("prefs", {"download.default_directory": DOWNLOAD_FOLDER})
    options.add_experimental_option('excludeSwitches', ['enable-logging'])
    if(FIREFOX):
        print("[WARNING] This is not currently supported for file matching, only listing, or proxy usage.\n          For use with Firefox some adaptations need to be made to the __init__ function.")
        driver = webdriver.Firefox()
    else:
        wdriver = os.getenv("WEBDRIVER")
        if(not wdriver):
            wdriver = "chromedriver.exe" if os.name == "nt" else "chromedriver"
        driver = webdriver.Chrome(wdriver, chrome_options=options) 
    driver.get(URL)
    return driver

def __office365userAction__(driver, user):
    driver.find_element_by_css_selector("[type=email]").send_keys(user)
    driver.find_elements_by_css_selector("[type=submit]")[0].click()

def __userAction__(driver, user):
    actions = ActionChains(driver)
    actions.send_keys(Keys.TAB)
    actions.send_keys(Keys.TAB)
    actions.send_keys(user)
    actions.send_keys(Keys.ENTER)
    return actions

def __confirmAction__(driver):
    driver.find_element_by_css_selector("[type=submit][value=Yes]").click()

def __redirectAction__(driver, user):
    driver.find_elements_by_css_selector("[type=submit]")[0].click()

def __passwordAction__(driver, password):
    actions = ActionChains(driver)
    actions.send_keys(password)
    actions.send_keys(Keys.ENTER)
    return actions

OFFICE365_PAGE = lambda driver : driver.current_url.lstrip("https://").startswith("login.microsoftonline.com/")

def __login__(driver, user, password):
    print("[+] Logging in.")
    if(not user):
        print("Missing user.\nPlease enter user:", end="")
        user = input()
    office365 = OFFICE365_PAGE(driver)
    if(office365):
        __office365userAction__(driver, user)
    else:
        __userAction__(driver, user).perform()
    time.sleep(4*DELAY)
    if(not office365 and OFFICE365_PAGE(driver)): # Redirected to Office 365 page; Click next again
        __redirectAction__(driver, user)          # TODO return this redirection information to main context so proper GetFS function can be selected.
        time.sleep(4*DELAY)
    if(not password):
        print("Missing password.\nPlease enter password or finish login manually and press enter to continue.\n> ", end="")
        password = input()
    __passwordAction__(driver, password).perform()
    time.sleep(4*DELAY)
    if("Please type in the code displayed on your authenticator app from your device" in driver.page_source):
        print("2FA detected.\nPlease finish login manually and press enter to continue.")
        input()
    if("Stay signed in?" in driver.page_source):
        __confirmAction__(driver)
        time.sleep(4*DELAY)
    print("[+] Login routine complete")
    


opened = []

DEBUG_LEVEL = -1
DELAY = 1
PERMANENT_DOWNLOAD_FOLDER = False

# Yes, this has XSS
GetFS = lambda driver,path="" : driver.execute_script("""
r = document.getElementById("appRoot")
l = r.querySelectorAll("[data-list-index]")
fs = {true:[], false: [], "path": "%PATH%"}
for(var i = 0; i < l.length; i++) { 
    t = l[i].getElementsByClassName("FileTypeIcon")[0]["title"]
    s = l[i].querySelector("[data-automation-key=size]").getElementsByTagName("span")[0].textContent
    lnk = l[i].querySelector("[role=link]")
    if(t === "Folder")
        fs[true] = fs[true].concat({"name": lnk.title, "link": lnk.href, "path": `${fs["path"]}/${lnk.title}`, fs: {true: [], false: [], "path": `${fs["path"]}/${lnk.title}`, "size": s}})
    else
        fs[false] = fs[false].concat({"name": lnk.title, "link": lnk.href, "path": `${fs["path"]}/${lnk.title}`, "size": s})
}
return fs
""".replace("%PATH%", path))

# BROKEN the link element doesn't actually have a link on it.
GetOffice365FS = lambda driver, path="" : driver.execute_script("""
r = document.getElementById("appRoot")
l = document.querySelectorAll("[data-automationid=ListCell]")
fs = {true:[], false: [], "path": "%PATH%"}
for(var i = 0; i < l.length; i++) { 
    t = l[i].querySelector("[data-automationid=DetailsRowCell]").getElementsByTagName("i")[0].ariaLabel
    s = l[i].querySelector("[data-automation-key=fileSizeColumn_716]").textContent
    lnk = l[i].querySelector("[role=link]")
    if(t === "Folder")
        fs[true] = fs[true].concat({"name": lnk.title, "link": lnk.href, "path": `${fs["path"]}/${lnk.title}`, fs: {true: [], false: [], "path": `${fs["path"]}/${lnk.title}`, "size": s.replace("items","")}})
    else
        fs[false] = fs[false].concat({"name": lnk.title, "link": lnk.href, "path": `${fs["path"]}/${lnk.title}`, "size": s})
}
return fs
""".replace("%PATH%", path))

def __shallow_crawl(driver, folder): # BROKEN by extension until GetOffice365FS is fixed, and BROKEN as it does not handle a different GetFS function (Just pass the GetFS function to be used as a parameter.)
    global DELAY
    driver.get(folder["link"])
    time.sleep(DELAY)
    try:
        return GetFS(driver, folder["path"])
    except:
        print(f"[!] Exception while crawling {folder['name']}")
        if(DEBUG_LEVEL > 0):
            print(traceback.format_exc())
        return {"true": [], "false": [], "path": folder["path"], "size": "-1 KB"}


def crawl(driver, fs, blacklist=[], max_depth=-1, depth=0): # BROKEN by extension until GetOffice365FS is fixed, and BROKEN as it does not handle a different GetFS function (Just pass the GetFS function to be used as a parameter.)
    global DEBUG_LEVEL
    if(depth > max_depth and max_depth > -1):
        return fs
    if(DEBUG_LEVEL > 1):
        print(f"[DEBUG] Current at depth: {depth}")
    if(fs["path"] in blacklist):
        if(DEBUG_LEVEL > 0):
            print(f"[DEBUG] Blacklisted item found \"{fs['path']}\". Skipping.")
        return fs
    #ret = {"true": [], "false": []}
    for folder in fs["true"]: # First list all folders and files
        folder["fs"] = __shallow_crawl(driver, folder)
    for folder in fs["true"]: # Then crawl them
        folder["fs"] = crawl(driver, folder["fs"], blacklist, max_depth, depth+1)
    return fs

def __waitFile(fl): # This will brick if download fails # FIXED (Will still break if browser just dies and leaves the crdownload file there or the connection hangs, though)
    fn = os.path.join(DOWNLOAD_FOLDER, fl["name"])
    while(not os.path.exists(fn)):
        for i in os.listdir(DOWNLOAD_FOLDER):
            if (i.endswith(".crdownload") or i.endswith("tmp")) and i != fn: # fn could end with .crdownload, or tmp
                continue
        if(not os.path.exists(fn)):
            raise IOError("Missing download file")
        time.sleep(DELAY)

def __getFilePreview(driver, fl): # BROKEN for Sharepoint. Completely different preview URLs
    driver.get(f"{fl['link']}&o=OneUp")
    time.sleep(DELAY)
    driver.execute_script("""return document.querySelector("[data-automationid=download]").click()""")
    time.sleep((15 + humanfriendly.parse_size(fl["size"])/1000000/2)*DELAY)
    __waitFile(fl)
    #time.sleep(DELAY/2) # TODO: Continue this developlment line, should be a lot more reliable with some tweaking
    #driver.get("chrome://downloads")
    #time.sleep(DELAY/5)
    #items = driver.execute_script("""return document.querySelector("downloads-manager").shadowRoot.querySelector("#downloadsList").getElementsByTagName("downloads-item")""")
    

def __matchSimpleFile(content, matchers):
    global DOWNLOAD_FOLDER
    ret = {}
    c = content.split(b"\n")
    matched = []
    for matcher in matchers:
        m = re.compile(matcher)
        i = 0
        for line in c:
            r = m.findall(line)
            if(len(r) > 0):
                if(i in matched): # remove duplicate lines if multiple matches
                    continue
                ret[i] = line.strip()
                matched.append(i)
            i += 1
    return ret

def __matchParagraphs(paragraphs, m, matched_paragraphs, ret, prefix=""):
    i = 0
    matched = matched_paragraphs
    for p in paragraphs:
            if(i in matched):
                i += 1
                continue
            r = m.findall(p.text)
            if(len(r) > 0):
                ret[i] = p.text.strip()
                matched.append(f"{prefix}{i}")
            i += 1
    return (ret, matched)

def __matchTable(table, m, matched_rows, ret, prefix=""):
    matched = matched_rows
    ir = 0
    for row in table.rows:
        rid = f"{prefix}{ir}"
        if(rid in matched):
            ir += 1
            continue
        line = "||".join([cell.text for cell in row.cells])
        r = m.findall(line)
        if(len(r) > 0):
            ret[rid] = line.strip()
            matched.append(rid)
        ir += 1
    return (ret, matched)
        

def __matchTables(tables, m, matched_rows, ret, prefix=""):
    it = 0
    matched = matched_rows
    for t in tables:
        ret, matched = __matchTable(t, m, matched, ret, f"{prefix}{it}:")
        it += 1
           
    return (ret, matched)


def __matchWordFile(content, matchers):
    ret = {}
    doc = docx.Document(content)
    matched_paragraphs = []
    matched_rows = []
    for matcher in matchers:
        m = re.compile(matcher)
        ret, matched_paragraphs = __matchParagraphs(doc.paragraphs, m, matched_paragraphs, ret)
        ret, matched_rows = __matchTables(doc.tables, m, matched_rows, ret)
    return ret

def __matchPowerPointFile(content, matchers):
    ret = {}
    doc = pptx.Presentation(content)
    matched_paragraphs = []
    matched_rows = []
    for matcher in matchers:
        i = 0
        m = re.compile(matcher)
        for slide in doc.slides:
            isp = 0
            for shape in slide.shapes:
                sid = f"{i}:{isp}"
                if(shape.has_text_frame):
                    l = 0
                    for line in shape.text.split("\n"):
                        sid = f"{i}:{isp}.{l}"
                        if(sid in matched_paragraphs):
                            l += 1
                            continue
                        r = m.findall(shape.text)
                        if(len(r) > 0):
                            ret[sid] = line
                            matched_paragraphs.append(sid) # actually, these are useless now; change it later
                        l += 1
                        continue
                    isp += 1
                if(shape.has_table):
                    ret, matched_rows = __matchTable(shape.table, m, matched_rows, ret, f"{i}:{isp}:")
                    isp += 1
                    continue
                if(shape.has_chart):
                    # I don't expect to find anything usefull here, if you want sensitive data, instead of passwords, keys, and similar, plz update this to handle it
                    isp += 1
                    continue
            i += 1
    return ret

def __matchExcelFile(content, matchers):
    ret = {}
    workbook = openpyxl.load_workbook(content)
    for matcher in matchers:
        i = 0
        m = re.compile(matcher)
        for worksheet in workbook:
            ir = 0
            for row in worksheet.iter_rows():
                rid = f"{i}:{ir}"
                if(rid in ret):
                    ir += 1
                    continue
                line = "||".join([(cell.value if cell.value != None else "") for cell in row])
                r = m.findall(line)
                if(len(r) > 0):
                    ret[rid] = line.strip()
                ir += 1
                continue
            i += 1
    return ret

def __matchPDFFile(content, matchers):
    ret = {}
    pdfr = PyPDF2.PdfFileReader(content)
    for matcher in matchers:
        i = 0
        m = re.compile(matcher)
        for page in pdfr.pages:
            l = 0
            lines = page.extractText().split("\n")
            for line in lines:
                lid = f"{i}:{l}"
                if(lid in ret):
                    l += 1
                    continue
                r = m.findall(line)
                if(len(r) > 0):
                    ret[lid] = line.strip()
                    l += 1
            i += 1
    return ret

def __matchMSGFile(content, matchers, recurse_attachments=True, extension_blacklist=[], skip_binary=True, prefix="", path=""):
    ret = {}
    try:
        msg = extract_msg.openMsg(content)
    except:
        if(type(content) is extract_msg.Message):
            msg = content
        else:
            raise
    for matcher in matchers:
        i = 0
        lines = msg.body.split("\n")
        m = re.compile(matcher)
        for line in lines:
            lid = f"{prefix}{i}"
            if(lid in ret):
                i += 1
                continue
            r = m.findall(line)
            if(len(r) > 0):
                ret[lid] = line.strip()
            i += 1
        sid = f"{prefix}s"
        if(sid in ret):
            continue
        r = m.findall(msg.subject)
        if(len(r) > 0):
            ret[sid] = msg.subject.strip()
    if(not recurse_attachments):
        return ret
    i = 0
    for attachment in msg.attachments:
        if(attachment.longFilename):
                aid = f"{prefix}{attachment.longFilename.replace(':', '%3A')}"
                patha = f"{path}:{attachment.longFilename.replace(':', '%3A')}"
        else:
            if(msg.filename):
                aid = f"{prefix}{msg.filename.replace(':', '%3A')}"
                patha = f"{path}:{msg.filename.replace(':', '%3A')}"
            else:
                aid = f"{prefix}{i}"
                patha = path
        if(attachment.type == "msg"):
            aid = f"{prefix}{attachment.data.filename.replace(':', '%3A')}"
            r = __matchMSGFile(attachment.data, matchers, recurse_attachments, extension_blacklist, skip_binary, f"{aid}:", patha)
        else:
            r = {aid: __matchSelector(attachment.longFilename, f"{path}:{aid}", attachment.data, matchers, extension_blacklist, skip_binary, recurse_attachments, f"{aid}:")}
        if(len(r) > 0):
            ret.update(r)
        i += 1
    return ret


def __matchSelector(fn, path, content, matchers, extension_blacklist=[], skip_binary=True, msg_recurse_attachments=True, prefix=""):
    for extension in extension_blacklist:
        if(fn.endswith(extension)):
            if(DEBUG_LEVEL > 1):
                print(f"[DEBUG] Blacklisted extension matched {extension} for {path}. Skipping.")
            return "Skipped"
    mg = magic.from_buffer(content)
    if(mg.startswith("Microsoft Word") or (mg.startswith("Zip archive data") and (fn.endswith(".doc") or fn.endswith(".docx")))):
        m = __matchWordFile(content, matchers)
    elif(mg.startswith("Microsoft PowerPoint") or (mg.startswith("Zip archive data") and (fn.endswith(".ppt") or fn.endswith("pptx")))):
        m = __matchPowerPointFile(content, matchers)
    elif(mg.startswith("Microsoft Excel") or (mg.startswith("Zip archive data") and (fn.endswith(".xls") or fn.endswith("xlsx")))):
        m = __matchExcelFile(content, matchers)
    elif(mg.startswith("PDF document")):
        m = __matchPDFFile(content, matchers)
    elif(mg.endswith("Microsoft Outlook Message")):
        m = __matchMSGFile(content, matchers, msg_recurse_attachments, extension_blacklist, skip_binary, prefix, path)
    else:
        if(mg == "data" and skip_binary):
            print(f"Binary file detected({path}). Skipping.")
            return "Skipped"
        matchers = [i.encode() for i in matchers] # Convert to binary form
        m = __matchSimpleFile(content, matchers)
    return m


### Returns {file name: matches}; matches is just the string "Skipped" if processing was skipped for whatever reason.
### No exceptions are treated other than UnicodeDecodeError byt the raw data matcher, all else should be treated by invoker.
def __matchFile(fl, matchers, extension_blacklist=[], skip_binary=True, msg_recurse_attachment=True):
    global PERMANENT_DOWNLOAD_FOLDER
    ret = {}
    content = ""
    fn = os.path.join(DOWNLOAD_FOLDER, fl["name"])
    if(DEBUG_LEVEL > 2):
        print(f"[DEBUG] Matching {fl['name']}")
    with open(fn, "rb") as fd:
        content = fd.read()
    m = __matchSelector(fl["name"], fl["path"], content, matchers, extension_blacklist, skip_binary, msg_recurse_attachment)
    if(len(m) > 0):
        ret = {fl["path"]: m}
    if(DEBUG_LEVEL > 2):
        print(f"[DEBUG] Finished matching file. Matches: {ret}.")
    if(PERMANENT_DOWNLOAD_FOLDER):
        i = 0
        while(True):
            try:
                s = "" if i == 0 else f"({i})"
                os.rename(fn, os.path.join(PERMANENT_DOWNLOAD_FOLDER, f"{fl['name']}{s}"))
            except FileExistsError:
                i += 1
    else:
        os.remove(fn)
    return ret


def match(driver, fs, matchers, blacklist=[], max_size=-1, max_depth=-1, depth=0, extension_blacklist=[], skip_binary=True, msg_recurse_attachment=True): # BROKEN by extension until __getFilePreview is fixed
    global DEBUG_LEVEL
    ret = {}
    if(depth > max_depth and max_depth > -1):
        if(DEBUG_LEVEL > 1):
            print(f"[DEBUG] Max depth reached at \"{fs['path']}\", depth: {depth}")
        return ret
    if(fs["path"] in blacklist):
        if(DEBUG_LEVEL > 0):
            print(f"[DEBUG] Blacklisted item found \"{fs['path']}\". Skipping.")
    for folder in fs["true"]:
        ret += match(driver, folder, matchers, blacklist, max_size, max_depth, depth+1)
    for fd in fs["false"]:
        try:
            if(max_size > -1):
                s = humanfriendly.parse_size(fd["size"])
                ms = humanfriendly.parse_size(max_size)
                if(s > ms):
                    if(DEBUG_LEVEL > 0):
                        print(f"[DEBUG] File \"{fd['path']}\" above max file size, size: {s}. Skipping.")
                    continue
            for extension in extension_blacklist:
                if(fd["path"].endswith(extension)):
                    if(DEBUG_LEVEL > 2):
                        print(f"[DEBUG] Blacklisted extension matched {extension} for {fd['path']}. Skipping download.")
                    continue
            __getFilePreview(driver, fd)
            ret.update(__matchFile(fd, matchers, extension_blacklist, skip_binary, msg_recurse_attachment))
        except:
            print(f"[ERROR] Failed to download {fd['path']}")
            if(DEBUG_LEVEL > 0):
                traceback.print_exc()
    return ret

def printFS(fs, indent=""):
    try:
        for folder in fs["true"]:
            print(f"{indent}+ {folder['name']}:")
            printFS(folder["fs"], f"{indent}|   ")
            print(f"{indent}\\")
        for fn in fs["false"]:
            print(f"{indent}{fn['name']}")
    except:
        for folder in fs["True"]:
            print(f"{indent}+ {folder['name']}:")
            printFS(folder["fs"], f"{indent}|   ")
            print(f"{indent}\\")
        for fn in fs["False"]:
            print(f"{indent}{fn['name']}")

def printM(m, indent=""):
    for path in m:
        print(f"{indent}+ {path}: ")
        i = 1
        if(type(m[path]) is str):
            print(f"{indent}\\    {m[path]}")
        for location in m[path]:
            s = "\\" if i == len(m[path]) else "|"
            print(f"{indent}{s}   {location}: {m[path][location]}")
            i += 1

#fs = {True: [], False: []}
def main():
    global DEBUG_LEVEL
    global DELAY
    global PERMANENT_DOWNLOAD_FOLDER
    global URL
    global PROXY
    global FIREFOX

    parser = argparse.ArgumentParser()
    parser.add_argument("--verbose", "-v", dest="verbose", action="count", default=0)
    parser.add_argument("--max-depth", "-d", metavar="N", dest="max_depth", type=int, default=-1)
    parser.add_argument("--blacklist", "-e", metavar="FILE_PATH", dest="exclude", action="append", default=[])
    parser.add_argument("--delay", "-t", metavar="T", dest="delay", type=float, default=1)
    parser.add_argument("--json", "-o", metavar="OUTPUT_FILE", dest="output", default=False)
    parser.add_argument("--matcher", "-m", metavar="REGEX_STRING", dest="matchers", action="append", default=[])
    parser.add_argument("--quiet", "-q", dest="quiet", action="store_true", default=False)
    parser.add_argument("--max-size", "-L", metavar="SIZE_LIMIT", dest="max_size", default=-1)
    parser.add_argument("--keep-files", "-k", metavar="PERMANENT_DOWNLOAD_FOLDER", dest="download_folder", default=False)
    parser.add_argument("--match-binary-files", "-A", dest="parse_binary", action="store_true", default=False)
    parser.add_argument("--disable-msg-recursion", "-f", dest="disable_msg_recursion", action="store_true", default=False)
    parser.add_argument("--extension-blacklist", "-x", metavar=".EXTENSION", dest="extension_blacklist", action="append", default=[])
    parser.add_argument("--load-fs", "-l", metavar="FS_JSON", dest="fs_json", default=False)
    parser.add_argument("--office-365", "-b", metavar="LOGIN_URL", dest="business_login", default=False)
    parser.add_argument("--user", "-u", metavar="USER", dest="user", default=False)
    parser.add_argument("--password", "-p", metavar="PASSWORD", dest="password", default=False)
    parser.add_argument("--proxy", "-P", metavar="PROXY_URI", dest="proxy", default=False)
    parser.add_argument("--firefox", "-F", dest="firefox", action="store_true", default=False)

    args = parser.parse_args()

    DEBUG_LEVEL = args.verbose
    DELAY = args.delay
    PERMANENT_DOWNLOAD_FOLDER = args.download_folder
    PROXY = args.proxy
    FIREFOX = args.firefox
    if(PERMANENT_DOWNLOAD_FOLDER):
        try:
            os.mkdir(PERMANENT_DOWNLOAD_FOLDER)
        except FileExistsError:
            pass

    if(args.quiet and not args.output):
        print("To use the --quiet/-q flag enable file output with --json/-o {OUTPUT_FILE}")
        return
    try:
        humanfriendly.parse_size(args.max_size)
    except:
        print(f"Invalid max-size: {args.max_size}.\nPlease enter a valid size indentifier. IE: 20KB, 20000, 1.5MB")
        return

    if(args.business_login):
        URL = args.business_login

    driver = __init__()
    time.sleep(4*DELAY)
    __login__(driver, args.user, args.password)

    if(args.fs_json):
        print(f"[+] Loading file system from {args.fs_json}")
        fs = json.load(args.fs_json)
    else:
        print(f"[+] Extracting file system root from OneDrive")
        fs = GetFS(driver)
        print(f"[+] Crawling folders")
        fs = crawl(driver, fs, args.exclude, args.max_depth)
        print(f"[+] Extraction complete")

    if(not args.quiet):
        printFS(fs)

    output = False
    if(args.output):
        output = args.output[:-5] if args.output.endswith(".json") else args.output
        outputfs = f"{output}_fs.json"
        with open(outputfs, "a+") as fd:
            json.dump(fs, fd)

    if(len(args.matchers) <= 0):
        return fs
    

    m = match(driver, fs, args.matchers, args.exclude, args.max_size, args.max_depth)
    if(not args.quiet):
        printM(m)

    if(output):
        outputm = f"{output}_matches.json"
        with open(outputm, "a+") as fd:
            json.dump(m, fd)


    return (fs, m)

if(__name__ == "__main__"):
    main()